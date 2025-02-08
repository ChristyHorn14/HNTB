import glog
import pandas as pd
from pptx import Presentation

from hntb.config_options import HNTBConfig


# Function to duplicate a slide
def duplicate_slide(prs, slide):
    slide_layout = prs.slide_layouts[0]
    new_slide = prs.slides.add_slide(slide_layout)
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            new_table = new_slide.shapes.add_table(
                rows=len(table.rows),
                cols=len(table.columns),
                left=shape.left,
                top=shape.top,
                width=shape.width,
                height=shape.height,
            ).table
            for r in range(len(table.rows)):
                for c in range(len(table.columns)):
                    new_table.cell(r, c).text = table.cell(r, c).text
    return new_slide


# Define a function to replace placeholders in table cells
def replace_placeholder_in_table(table, mapping):
    for row in table.rows:
        for cell in row.cells:
            for placeholder, value in mapping.items():
                if placeholder in cell.text:
                    cell.text = cell.text.replace(placeholder, value)


def generate_ppt(cfg: HNTBConfig):
    # Load the PowerPoint template
    template_file = cfg.template_directory / cfg.ppt_template_filename
    glog.debug(f"=> template_file={template_file}")
    prs = Presentation(template_file)

    # Load the Excel file
    input_file = cfg.active_tumor_board_file
    output_file = cfg.output_directory / cfg.ppt_filename
    df = pd.read_excel(input_file, sheet_name="Master Linked")

    glog.debug(f"=> input_file={input_file}")
    glog.debug(f"=> output_file={output_file}")

    # Find the template slide with placeholders
    template_slide = None
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if "{" in cell.text and "}" in cell.text:
                            template_slide = slide
                            break
                if template_slide:
                    break
        if template_slide:
            break

    # Iterate through each row in the dataframe
    for index, row in df.iterrows():
        # Duplicate the template slide for each row in the dataframe
        slide = duplicate_slide(prs, template_slide)

        # Create a mapping of placeholders to actual values
        mapping = {
            "{Initials}": str(row["Initials"]),
            "{Demographics}": str(row["Demographics"]),
            "{Diagnosis}": str(row["Diagnosis"]),
            "{Attending}": str(row["Attending"]),
        }

        # Replace placeholders in tables with actual values from the dataframe
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                replace_placeholder_in_table(table, mapping)

    # Save the modified PowerPoint presentation
    prs.save(output_file)
    glog.info(f"=> Document saved successfully to {output_file}")
